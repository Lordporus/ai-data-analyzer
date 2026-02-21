"""
PowerPoint generation module — Create a professional slide deck
from pipeline results using python-pptx.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config.settings import BRAND_NAME


# Brand colors
BRAND_PURPLE = RGBColor(0x6C, 0x63, 0xFF)
DARK_BG = RGBColor(0x0D, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1B, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x8B, 0x94, 0x9E)


def generate_pptx(
    pipeline_result,
    output_dir: str | Path,
) -> str:
    """Generate a PowerPoint deck from pipeline results.

    Args:
        pipeline_result: PipelineResult from the orchestrator.
        output_dir: Directory to save the PPTX file.

    Returns:
        Path to the generated .pptx file.
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "report.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_bg(slide, DARK_BG)

    _add_text(
        slide, BRAND_NAME, 1.5, 2.0, 10.0, 1.5,
        font_size=44, bold=True, color=BRAND_PURPLE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, "Automated Data Analysis Report", 1.5, 3.5, 10.0, 1.0,
        font_size=24, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        f"Generated {datetime.now().strftime('%B %d, %Y')}",
        1.5, 4.8, 10.0, 0.6,
        font_size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )

    # ── Slide 2: Executive Summary ───────────────────────────────────
    ingestion = pipeline_result.ingestion
    if ingestion:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, DARK_BG)
        _add_text(
            slide, "Executive Summary", 0.8, 0.5, 11.0, 1.0,
            font_size=32, bold=True, color=BRAND_PURPLE,
        )
        summary_text = (
            f"Dataset: {ingestion.row_count:,} rows x {ingestion.col_count} columns\n"
            f"Duplicates found: {ingestion.duplicate_count:,}\n"
            f"Missing values: {sum(ingestion.missing_value_report.values()):,}\n"
            f"File size: {ingestion.file_size_bytes / 1024:.1f} KB"
        )
        _add_text(
            slide, summary_text, 0.8, 1.8, 11.0, 3.5,
            font_size=18, color=WHITE,
        )

    # ── Slide 3: Key Metrics ─────────────────────────────────────────
    if pipeline_result.insight:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, DARK_BG)
        _add_text(
            slide, "Key Metrics", 0.8, 0.5, 11.0, 1.0,
            font_size=32, bold=True, color=BRAND_PURPLE,
        )
        kpis = pipeline_result.insight.kpi_list[:10]
        kpi_text = "\n".join(f"• {k.name}: {k.value}" for k in kpis)
        _add_text(
            slide, kpi_text, 0.8, 1.8, 11.0, 4.5,
            font_size=16, color=WHITE,
        )

    # ── Slide 4: Recommendations ─────────────────────────────────────
    if pipeline_result.insight and pipeline_result.insight.business_recommendations:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, DARK_BG)
        _add_text(
            slide, "Recommendations", 0.8, 0.5, 11.0, 1.0,
            font_size=32, bold=True, color=BRAND_PURPLE,
        )
        recs = pipeline_result.insight.business_recommendations[:6]
        # Strip emoji for PPTX compatibility
        recs_text = "\n\n".join(
            r.encode("ascii", "ignore").decode("ascii").strip()
            for r in recs if r.encode("ascii", "ignore").decode("ascii").strip()
        )
        _add_text(
            slide, recs_text, 0.8, 1.8, 11.0, 4.5,
            font_size=14, color=WHITE,
        )

    # ── Slide 5: Thank You ───────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_text(
        slide, "Thank You", 1.5, 2.5, 10.0, 1.5,
        font_size=44, bold=True, color=BRAND_PURPLE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, f"Report by {BRAND_NAME}", 1.5, 4.5, 10.0, 0.8,
        font_size=16, color=MUTED, align=PP_ALIGN.CENTER,
    )

    prs.save(str(pptx_path))
    return str(pptx_path)


# ── Helpers ──────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide, text: str,
    left: float, top: float, width: float, height: float,
    font_size: int = 18, bold: bool = False,
    color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
) -> None:
    """Add a text box to the slide."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
